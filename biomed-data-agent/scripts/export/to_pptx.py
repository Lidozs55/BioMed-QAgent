"""PPT 报告生成脚本（BioMed QAgent Stage 6，可选）。

输入：cleaned DataRecord JSON + 可选 lineage.json + 可选分析结果目录
输出：report.pptx（PowerPoint 格式，中文，不超过 20 页）

复用 to_report.py 的 build_report() 构造报告结构，再用 python-pptx 渲染为 PPT。
python-pptx 不可用时输出错误并建议安装该依赖或改用调度器自带的文档处理能力。

页面规划（对应赛题"不超过 20 页"约束）：
  - 第 1 页：封面（标题 + 任务 ID + 时间）
  - 第 2 页：执行摘要
  - 第 3 页：数据源统计
  - 第 4 页：字段映射
  - 第 5 页：质量审核
  - 第 6 页：数据溯源
  - 第 7-12 页：分析结果（最多 6 页，每页一个分析）
  - 第 13 页：输出文件清单
  - 第 14 页：附录字段映射表（如有）
  超出 20 页时自动截断并在末页提示。

接口：
    python scripts/export/to_pptx.py --input cleaned.json \
        --lineage lineage.json --analysis-dir results/ \
        --out report.pptx --task-id T1

成功输出：{"status":"ok","output":"report.pptx","rows":N,"slides":M}
失败输出：{"status":"error","message":"..."}
"""
import os
import sys

_HERE = os.path.dirname(os.path.abspath(__file__))
if _HERE not in sys.path:
    sys.path.insert(0, _HERE)
from _base import emit_error, emit_ok, load_lineage, load_records, log_stderr  # noqa: E402
from to_report import build_report  # noqa: E402

MAX_SLIDES = 20


def _add_title_slide(prs, title_text, subtitle_text):
    """添加封面页。"""
    from pptx.util import Inches, Pt
    slide_layout = prs.slide_layouts[0]  # Title Slide
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    if title.has_text_frame:
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(36)
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = subtitle_text
        if subtitle.has_text_frame:
            for p in subtitle.text_frame.paragraphs:
                for r in p.runs:
                    r.font.size = Pt(16)


def _add_content_slide(prs, title_text, blocks):
    """添加内容页（标题 + 文本/表格 blocks）。"""
    from pptx.util import Inches, Pt
    from pptx.dml.color import RGBColor
    slide_layout = prs.slide_layouts[1]  # Title and Content
    slide = prs.slides.add_slide(slide_layout)
    title = slide.shapes.title
    title.text = title_text
    if title.has_text_frame:
        for p in title.text_frame.paragraphs:
            for r in p.runs:
                r.font.size = Pt(28)
                r.font.color.rgb = RGBColor(0x1A, 0x4D, 0x8F)
    body = slide.placeholders[1]
    tf = body.text_frame
    tf.word_wrap = True
    first = True
    for b in blocks:
        btype = b.get("type")
        if btype == "para":
            p = tf.paragraphs[0] if first else tf.add_paragraph()
            p.text = b.get("text", "")
            for r in p.runs:
                r.font.size = Pt(14)
            first = False
        elif btype == "bullets":
            for it in b.get("items", []):
                p = tf.paragraphs[0] if first else tf.add_paragraph()
                p.text = str(it)
                p.level = 0
                for r in p.runs:
                    r.font.size = Pt(13)
                first = False
        elif btype == "table":
            headers = b.get("headers", [])
            rows = b.get("rows", [])
            n_rows, n_cols = len(rows) + 1, len(headers)
            if n_rows < 2 or n_cols < 1:
                continue
            left, top, width, height = Inches(0.5), Inches(3.5), Inches(9), Inches(3.5)
            tbl_shape = slide.shapes.add_table(n_rows, n_cols, left, top, width, height)
            tbl = tbl_shape.table
            for ci, h in enumerate(headers):
                cell = tbl.cell(0, ci)
                cell.text = str(h)
                for p in cell.text_frame.paragraphs:
                    for r in p.runs:
                        r.font.size = Pt(11)
                        r.font.bold = True
            for ri, row in enumerate(rows, start=1):
                for ci, val in enumerate(row):
                    if ci >= n_cols:
                        break
                    cell = tbl.cell(ri, ci)
                    cell.text = str(val)
                    for p in cell.text_frame.paragraphs:
                        for r in p.runs:
                            r.font.size = Pt(10)


def render_pptx(sections, output_path, task_id=""):
    """用 python-pptx 把 sections 渲染为 PPT 文档。"""
    from pptx import Presentation
    from pptx.util import Inches, Pt
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    if not sections:
        _add_title_slide(prs, "生物医学数据整合报告", f"任务 ID: {task_id or '(未指定)'}")
    else:
        # 第 1 页：封面（用第一个 section 的标题作为主标题）
        head = sections[0]
        head_bullets = []
        for b in head["blocks"]:
            if b.get("type") == "bullets":
                head_bullets.extend(b.get("items", []))
        _add_title_slide(prs, head["title"], "\n".join(head_bullets[:4]))

    # 后续页：每个 section 一页（控制在 MAX_SLIDES-1 页内）
    remaining = MAX_SLIDES - 1
    for sec in sections[1:]:
        if remaining <= 0:
            log_stderr(f"已达 {MAX_SLIDES} 页上限，剩余 section 已截断")
            # 末页提示
            _add_content_slide(prs, "（内容截断）",
                               [{"type": "para",
                                 "text": f"本 PPT 已达 {MAX_SLIDES} 页上限，"
                                         "完整内容请见 report.md / report.docx / report.pdf。"}])
            break
        # 控制单页 block 数量避免溢出
        blocks = sec["blocks"]
        if len(blocks) > 8:
            blocks = blocks[:8]
            _add_content_slide(prs, sec["title"], blocks +
                               [{"type": "para", "text": "…（已截断，详见完整报告）"}])
        else:
            _add_content_slide(prs, sec["title"], blocks)
        remaining -= 1

    prs.save(output_path)
    return len(prs.slides)


def main():
    import argparse
    parser = argparse.ArgumentParser(prog="to_pptx", description="生成生物医学数据整合 PPT 报告（不超过 20 页）")
    parser.add_argument("--input", required=True, help="输入 cleaned DataRecord JSON")
    parser.add_argument("--lineage", default=None, help="lineage.json 路径（可选）")
    parser.add_argument("--analysis-dir", default=None, help="分析结果目录（可选）")
    parser.add_argument("--out", required=True, help="输出 report.pptx 路径")
    parser.add_argument("--task-id", default="", help="任务 ID")
    args = parser.parse_args()
    try:
        try:
            import pptx  # noqa: F401
        except ImportError:
            emit_error("python-pptx 未安装。建议：(1) 沙箱中执行 pip install python-pptx；"
                       "或 (2) 改用调度器自带的文档处理能力生成 PPT 报告。")
        records = load_records(args.input)
        log_stderr(f"加载 {len(records)} 条记录")
        sections = build_report(records, load_lineage(args.lineage), args.task_id,
                                args.analysis_dir, args.out, input_path=args.input)
        os.makedirs(os.path.dirname(os.path.abspath(args.out)) or ".", exist_ok=True)
        n_slides = render_pptx(sections, args.out, task_id=args.task_id)
        log_stderr(f"PPT 报告已写入 {args.out}（{n_slides} 页）")
        emit_ok(args.out, rows=len(records), sections=len(sections), slides=n_slides)
    except Exception as e:
        emit_error(f"PPT 报告生成失败: {e}")


if __name__ == "__main__":
    main()
